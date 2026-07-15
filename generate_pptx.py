import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme colors
    COLOR_BG = RGBColor(15, 23, 42)      # Slate 900
    COLOR_TITLE = RGBColor(56, 189, 248)  # Cyan 400
    COLOR_TEXT = RGBColor(241, 245, 249)   # Slate 50
    COLOR_MUTED = RGBColor(148, 163, 184) # Slate 400

    def add_slide_with_bg():
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return slide

    def add_title(slide, text):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = COLOR_TITLE
        
        # Add accent line under title
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.75), Inches(1.4), Inches(11.83), Inches(0.04)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = COLOR_TITLE
        accent.line.fill.background()

    def add_content(slide, bullets, left=Inches(0.75), top=Inches(1.8), width=Inches(11.83), height=Inches(5.0)):
        content_box = slide.shapes.add_textbox(left, top, width, height)
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            level = 0
            text = bullet
            if bullet.startswith("    - "):
                level = 2
                text = bullet[6:]
            elif bullet.startswith("  - "):
                level = 1
                text = bullet[4:]
            elif bullet.startswith("- "):
                level = 0
                text = bullet[2:]
                
            p.text = text
            p.level = level
            p.space_after = Pt(12)
            p.font.name = 'Segoe UI'
            p.font.size = Pt(18 - level * 2)
            p.font.bold = (level == 0) # bold for main bullets
            p.font.color.rgb = COLOR_TEXT if level == 0 else COLOR_MUTED

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG
    bg1.line.fill.background()
    
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(1.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "BudgetBrake"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = 'Segoe UI'
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TITLE
    
    subtitle_box = slide1.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11.33), Inches(1.0))
    tf_sub = subtitle_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Shop till the Budget Stops"
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.name = 'Segoe UI'
    p_sub.font.size = Pt(22)
    p_sub.font.italic = True
    p_sub.font.color.rgb = COLOR_TEXT
    
    info_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.33), Inches(2.0))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    info_lines = [
        "A Native, Secure C-Win32 Desktop GUI Personal Finance & Shopping Manager",
        "Built for C Programming Lab | Project Based Learning (PBL)",
        "Presented by: Samarth Kamble"
    ]
    for i, line in enumerate(info_lines):
        if i == 0:
            p_info = tf_info.paragraphs[0]
        else:
            p_info = tf_info.add_paragraph()
        p_info.text = line
        p_info.alignment = PP_ALIGN.CENTER
        p_info.font.name = 'Segoe UI'
        p_info.font.size = Pt(14)
        p_info.font.color.rgb = COLOR_MUTED
        if i == len(info_lines)-1:
            p_info.font.color.rgb = COLOR_TITLE
            p_info.font.bold = True

    # Slide 2: Introduction to C
    slide2 = add_slide_with_bg()
    add_title(slide2, "Introduction to C")
    bullets2 = [
        "- History and Origin:",
        "  - Developed by Dennis Ritchie at Bell Laboratories in 1972.",
        "  - Foundation for many modern systems and compiled languages.",
        "- Core Characteristics:",
        "  - Statically-typed, compiled, and highly structured procedural language.",
        "  - Offers low-level hardware control, direct pointer manipulation, and memory management.",
        "- Relevance to GUI Applications:",
        "  - Allows direct interfacing with OS-specific libraries like the Windows API.",
        "  - Enables developers to write lightweight, standalone desktop applications with zero runtime dependencies and an extremely low memory footprint (<5 MB RAM)."
    ]
    add_content(slide2, bullets2)

    # Slide 3: Problem Statement
    slide3 = add_slide_with_bg()
    add_title(slide3, "Problem Statement")
    bullets3 = [
        "- Expense Tracking Issues:",
        "  - Consumers frequently struggle to monitor cumulative costs in real-time when buying items across multiple categories and shops.",
        "- Post-Facto Budget Tracking:",
        "  - Most budgeting tools only record expenses after purchases have occurred, failing to prevent overspending at the exact moment of sale.",
        "- Bulky Development Frameworks:",
        "  - Modern desktop applications often rely on heavy browser runtimes (like Electron or JS engines) that consume hundreds of megabytes of RAM.",
        "- The Project Goal:",
        "  - Design a secure, lightweight, and offline Windows GUI application in C that integrates a shopping cart simulator with real-time budget blocking."
    ]
    add_content(slide3, bullets3)

    # Slide 4: Objectives
    slide4 = add_slide_with_bg()
    add_title(slide4, "Objectives")
    bullets4 = [
        "- Real-Time Budget Enforcement:",
        "  - Intercept and block checkouts automatically if the transaction total exceeds the user's remaining monthly allowance.",
        "- Local Data Persistence:",
        "  - Save and load user profiles, account credentials, category catalogs, and shopping histories to flat text files.",
        "- Secure Session Management:",
        "  - Implement user registration, passwords, and data isolation between accounts (e.g., samarth, krithika).",
        "- Native Graphics Visualization:",
        "  - Render a visual expense graph (Cost vs. Time) directly inside the GUI using Windows GDI without third-party chart libraries.",
        "- Optimized Software Architecture:",
        "  - Build a standalone .exe with minimal footprint, high performance, and rapid startup time."
    ]
    add_content(slide4, bullets4)

    # Slide 5: Important Features of C Used
    slide5 = add_slide_with_bg()
    add_title(slide5, "Important Features of C Used in Project")
    bullets5 = [
        "- Structs (Structures):",
        "  - Grouped related properties to manage entities, such as the User profile struct (username, password, budget, log array) and Shop inventory structures.",
        "- File I/O (fopen, fprintf, fscanf):",
        "  - Acted as the local database system to read and write user settings, transaction lists, and catalog details.",
        "- Pointers & Dynamic Referencing:",
        "  - Handled memory addresses to interact with the event handler (HWND, HDC) and pass data records between parsing modules without copying overhead.",
        "- Multidimensional Arrays:",
        "  - Maintained the cart index, active shop categories, and matrixes used to calculate analytics points for graph coordinates.",
        "- String Manipulation & Time Functions:",
        "  - Used string.h to parse CSV structures safely and time.h to log transaction logs with precise timestamps."
    ]
    add_content(slide5, bullets5)

    # Slide 6: Results
    slide6 = add_slide_with_bg()
    add_title(slide6, "Result")
    bullets6 = [
        "- Fully Functional Custom Win32 Desktop GUI:",
        "  - Implemented a clean Dark-Themed interface featuring owner-drawn controls, listboxes, and category selectors.",
        "- Proactive Budget Safeguards:",
        "  - Shows remaining budget updating dynamically. Successfully prevents checkouts if limits are breached, prompting warning boxes.",
        "- Integrated Analytics Dashboard:",
        "  - Plotted spending history graphs dynamically on a canvas coordinate system, tracking user spendings across different dates.",
        "- Dynamic Shop Customization:",
        "  - Inventories are read dynamically from files (e.g., Apple,120 in data/Groceries/Shop1.txt). Changes in inventory reflect immediately in the application."
    ]
    add_content(slide6, bullets6)

    # Slide 7: Conclusion
    slide7 = add_slide_with_bg()
    add_title(slide7, "Conclusion")
    bullets7 = [
        "- Successful Integration of Shopping & Budgeting:",
        "  - BudgetBrake successfully delivers a dual shopping experience and budget logger in a cohesive application.",
        "- Demonstration of C's Performance:",
        "  - Proves that native C using standard OS APIs runs faster, uses fewer system resources (<5 MB RAM), and compiles into a tiny standalone binary.",
        "- Key Educational Learning Outcomes:",
        "  - Developed advanced capabilities in Win32 event-driven programming, file-based data structures, standard library usage, and native UI graphics rendering."
    ]
    add_content(slide7, bullets7)

    # Slide 8: Future Scope and Work
    slide8 = add_slide_with_bg()
    add_title(slide8, "Future Scope and Work")
    bullets8 = [
        "- Relational Database Integration:",
        "  - Migrate from flat text files to a local database like SQLite for standard query optimization and scalability.",
        "- Enhanced Cryptographic Protocols:",
        "  - Implement password hashing algorithms (e.g. SHA-256) instead of saving text in plain files.",
        "- In-App Inventory Editor:",
        "  - Build a graphical administrator panel to allow merchants or users to edit price sheets and items directly within the GUI.",
        "- Cross-Platform Compilation:",
        "  - Port the Win32 graphical framework to cross-platform widgets (like Qt, GTK, or wxWidgets) to allow native operation on macOS and Linux."
    ]
    add_content(slide8, bullets8)

    output_path = "BudgetBrake_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully as '{output_path}'")

if __name__ == "__main__":
    create_presentation()
